from exporters.base import Exporter
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from utils.logger import get_logger
from utils.exceptions import exception_handler
from utils.helpers import split_by_h2, extract_text_blocks

logger = get_logger("pdf_exporter")

class PptExporter(Exporter):
    """
    專門處理 PPT (.pptx) 檔案匯出的類別。
    """

    @exception_handler
    def export(self, file_path: str):
        logger.info(f"Exporting PPT to {file_path}")
        prs = Presentation()

        # 封面
        cover_slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide Layout
        cover_slide.shapes.title.text = self.title
        cover_slide.placeholders[1].text = f"{self.period}\n{self.extra}"
        cover_slide.shapes.title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 內容
        for section, items in self.content.items():
            # 章節首頁
            section_slide = prs.slides.add_slide(prs.slide_layouts[5])
            section_slide.shapes.title.text = section

            for subtitle, body in items.items():
                for sub_title, content in split_by_h2(body) or [(subtitle, body)]:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = sub_title.strip()
                    tf = slide.shapes.placeholders[1].text_frame
                    tf.clear()

                    for para in extract_text_blocks(content):
                        p = tf.add_paragraph()
                        if para.startswith("• "):
                            p.text = para[2:].strip()
                            p.level = 1  # 條列縮排
                        else:
                            p.text = para.strip()

        prs.save(file_path)
        logger.info("PPT export completed")